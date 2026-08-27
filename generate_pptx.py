import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: 'python-pptx' library is required.")
    print("Please install it using: pip install python-pptx")
    sys.exit(1)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Palette (Matching MovieHub Website)
    BG_COLOR = RGBColor(13, 13, 16)       # #0D0D10 (Pitch Dark)
    CARD_BG = RGBColor(20, 20, 28)        # #14141C (Dark Card Surface)
    RED_ACCENT = RGBColor(229, 9, 20)     # #E50914 (Cinema Crimson Red)
    WHITE_TEXT = RGBColor(255, 255, 255)   # #FFFFFF (Crisp White)
    GRAY_TEXT = RGBColor(161, 161, 170)   # #A1A1AA (Muted Text)
    GOLD_ACCENT = RGBColor(245, 158, 11)  # #F59E0B (Amber Gold)
    BORDER_RED = RGBColor(180, 20, 30)

    slides_data = [
        {
            "slide_num": 1,
            "title": "MovieHub Platform",
            "subtitle": "Modern Full-Stack Cinema & Streaming Web Application",
            "category": "PROJECT PRESENTATION",
            "bullets": [
                "✦ Built with Angular 19, Node.js, Express & MongoDB",
                "✦ RESTful API Architecture with JWT Authentication",
                "✦ Admin Content Dashboard & Interactive Trailer Player",
                "✦ Cinema Dark Mode UI with Dynamic Filtering"
            ],
            "notes": "Good morning/afternoon Doctor. Today I am presenting my project 'MovieHub', a full-stack movie discovery and streaming web platform built with Angular, Node.js, Express, and MongoDB."
        },
        {
            "slide_num": 2,
            "title": "Project Vision & Executive Summary",
            "subtitle": "Addressing Modern Media Discovery Challenges",
            "category": "PROBLEM STATEMENT & OBJECTIVES",
            "bullets": [
                "• Challenge: Traditional movie portals suffer from static layouts, lack real-time management, and offer poor mobile UX.",
                "• Core Solution: MovieHub delivers a high-performance single-page app with instant dynamic content search.",
                "• User Engagement: Embedded HD trailer modal, rating reviews system, and custom favorites watchlist.",
                "• Admin Autonomy: Full administrative dashboard for real-time CRUD management of movie metadata."
            ],
            "notes": "In this slide, we cover the project vision. Traditional platforms often lack real-time updates and interactive media elements. MovieHub solves this by offering instant filtering, embedded YouTube trailers, and full admin autonomy."
        },
        {
            "slide_num": 3,
            "title": "System Architecture & Tech Stack",
            "subtitle": "Decoupled Full-Stack Architecture",
            "category": "TECHNOLOGY STACK",
            "bullets": [
                "• Frontend Layer: Angular 19, RxJS, TailwindCSS, Component-based Architecture.",
                "• Backend API: Node.js, Express.js Server, Modular Controllers & Routes.",
                "• Database Layer: MongoDB with Mongoose ODM (NoSQL Schema Design).",
                "• Security & Storage: JSON Web Tokens (JWT), Bcrypt Password Hashing, Multer File Uploads."
            ],
            "notes": "Here is the architectural overview. On the client side, we use Angular with standalone components and RxJS. The backend is an Express API connecting to MongoDB via Mongoose. Authentication is secured using JWT."
        },
        {
            "slide_num": 4,
            "title": "Key Features & User Experience",
            "subtitle": "Seamless Interaction & Cinematic Aesthetics",
            "category": "USER FEATURES",
            "bullets": [
                "• Dynamic Hero Banner: Rotates featured movies with high-resolution backdrops.",
                "• Category Browsing: Featured, Latest Releases, Top Rated, and Genre filtering.",
                "• Interactive Trailer Player: Embedded video playback overlay with backdrop glassmorphism.",
                "• Profile & Favorites: Personalized watchlist and profile image upload persistence."
            ],
            "notes": "Moving on to user features: MovieHub offers an immersive Hero banner with movie slider controls, interactive trailer modal, category breakdown, and personal user watchlists."
        },
        {
            "slide_num": 5,
            "title": "Database Schema & Entity Models",
            "subtitle": "Structured MongoDB Collections",
            "category": "DATABASE DESIGN",
            "bullets": [
                "• Users Collection: name, email, password (hashed), role (admin/user), profileImage.",
                "• Movies Collection: title, overview, releaseYear, rating, genres[], poster, backdrop, trailerUrl, cast[].",
                "• Favorites Collection: userId (Ref: User), movieId (Ref: Movie), createdAt timestamp.",
                "• Reviews Collection: userId, movieId, rating score (1-5 stars), text comment."
            ],
            "notes": "This slide details our MongoDB data modeling. We have four core schemas: Users, Movies, Favorites, and Reviews, with relational references establishing clean data integrity."
        },
        {
            "slide_num": 6,
            "title": "Backend REST API Architecture",
            "subtitle": "Secure & Scalable Endpoints",
            "category": "BACKEND CONTROLLERS",
            "bullets": [
                "• Auth Endpoints: POST /api/auth/register, POST /api/auth/login, PUT /api/auth/updateimage.",
                "• Movie Endpoints: GET /api/movies (filtering & search), GET /api/movies/:id, POST/PUT/DELETE.",
                "• Favorites Endpoints: GET /api/favorites, POST /api/favorites, DELETE /api/favorites/:id.",
                "• Middleware Protection: JWT auth middleware verifying bearer tokens for protected routes."
            ],
            "notes": "On the backend side, we designed RESTful controllers. Protected operations like uploading avatar pictures or modifying movie records pass through our JWT authentication middleware."
        },
        {
            "slide_num": 7,
            "title": "Frontend Modular Architecture",
            "subtitle": "Angular Components & Services",
            "category": "FRONTEND DESIGN",
            "bullets": [
                "• Reusable Components: Navbar, MovieCard, TrailerModal, SimilarMovies, Footer.",
                "• Core Pages: Home, Movies Catalog, Movie Details, Favorites, Profile, Admin Dashboard.",
                "• Centralized Services: MovieService (HTTP requests), UserService, AuthService.",
                "• Route Protection: Angular Guards restricting access to user and admin routes."
            ],
            "notes": "Our Angular frontend follows strict modular design. Components like MovieCard and TrailerModal are reusable across views. Services handle data streams using RxJS Observables."
        },
        {
            "slide_num": 8,
            "title": "Admin Management Dashboard",
            "subtitle": "Empowering Platform Administrators",
            "category": "ADMIN DASHBOARD",
            "bullets": [
                "• Dynamic CRUD Control: Add new releases, update movie details, or remove obsolete titles.",
                "• Genre Multi-Selection: Interactive badge selector for genre classification.",
                "• Trailer URL Binding: Seamless YouTube link integration for instant video rendering.",
                "• Role Guards: Server-side & client-side verification ensuring admin-only access."
            ],
            "notes": "The Admin Dashboard gives authorized administrators complete authority over the movie repository, including adding trailers, setting ratings, and editing movie metadata."
        },
        {
            "slide_num": 9,
            "title": "UI Design System & Aesthetics",
            "subtitle": "Cinematic Crimson Theme",
            "category": "DESIGN SYSTEM",
            "bullets": [
                "• Primary Dark Palette: Deep Charcoal (#0D0D10) background reducing eye strain.",
                "• Accent Highlights: Electric Red (#E50914) primary CTAs and Amber Gold (#F59E0B) ratings.",
                "• Glassmorphism: Modern translucent card overlays with subtle borders.",
                "• Micro-Animations: Hover scale transitions and spinning loading indicators."
            ],
            "notes": "Our UI design system uses a sleek dark cinematic color palette matching industry streaming platforms. We incorporated crimson highlights, glassmorphic cards, and micro-interactions."
        },
        {
            "slide_num": 10,
            "title": "Conclusion & Future Roadmap",
            "subtitle": "Project Deliverables & Expansion",
            "category": "SUMMARY & FUTURE SCOPE",
            "bullets": [
                "• Deliverables Met: Complete full-stack web application ready for production deployment.",
                "• Future Feature 1: Real-time HLS Video Streaming Server integration.",
                "• Future Feature 2: AI-driven Personalized Recommendation Engine based on user viewing history.",
                "• Thank You: Thank you Doctor for your guidance! Ready for Q&A."
            ],
            "notes": "To summarize, MovieHub is a robust, full-stack application. In the future, we plan to add HLS streaming and AI recommendations. Thank you Doctor, I'd be happy to answer any questions!"
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Slide Background Fill
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR
        bg_shape.line.fill.background()

        # Top Red Accent Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(0.4), Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RED_ACCENT
        bar.line.fill.background()

        # Header Box (Category & Title)
        header_box = slide.shapes.add_textbox(Inches(1.4), Inches(0.4), Inches(11), Inches(1.3))
        tf = header_box.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = data["category"]
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = RED_ACCENT

        p1 = tf.add_paragraph()
        p1.text = data["title"]
        p1.font.size = Pt(26)
        p1.font.bold = True
        p1.font.color.rgb = WHITE_TEXT

        p2 = tf.add_paragraph()
        p2.text = data["subtitle"]
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY_TEXT

        # Main Content Card Shape (Dark Glass Box)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_RED
        card.line.width = Pt(1.5)

        # Card Content Box
        content_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.5))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        for i, bullet in enumerate(data["bullets"]):
            p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE_TEXT if not bullet.startswith("✦") else GOLD_ACCENT
            p.space_after = Pt(16)

        # Slide Number Badge
        num_box = slide.shapes.add_textbox(Inches(11.8), Inches(0.4), Inches(0.8), Inches(0.5))
        np = num_box.text_frame.paragraphs[0]
        np.text = f"{data['slide_num']}/10"
        np.font.size = Pt(14)
        np.font.bold = True
        np.font.color.rgb = RED_ACCENT
        np.alignment = PP_ALIGN.RIGHT

        # Presenter Notes
        slide.notes_slide.notes_text_frame.text = data["notes"]

    output_path = "MovieHub_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
